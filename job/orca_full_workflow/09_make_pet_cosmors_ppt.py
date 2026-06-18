#!/usr/bin/env python3
"""Create a PPT summary for the PET/openCOSMO-RS workflow."""

from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
import csv


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "PET_openCOSMO_RS_workflow_summary.pptx"


COLORS = {
    "ink": RGBColor(31, 42, 55),
    "muted": RGBColor(91, 105, 122),
    "blue": RGBColor(37, 99, 235),
    "green": RGBColor(5, 150, 105),
    "red": RGBColor(190, 18, 60),
    "amber": RGBColor(180, 83, 9),
    "light": RGBColor(243, 246, 250),
    "line": RGBColor(205, 213, 223),
    "white": RGBColor(255, 255, 255),
}


def add_textbox(slide, left, top, width, height, text, size=18, bold=False, color=None, align=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    if align:
        p.alignment = align
    run = p.runs[0]
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or COLORS["ink"]
    return box


def add_title(slide, title, subtitle=None):
    add_textbox(slide, Inches(0.55), Inches(0.28), Inches(12.2), Inches(0.52), title, 24, True)
    if subtitle:
        add_textbox(slide, Inches(0.58), Inches(0.84), Inches(12.0), Inches(0.32), subtitle, 10, False, COLORS["muted"])


def add_footer(slide, text):
    add_textbox(slide, Inches(0.55), Inches(7.15), Inches(12.2), Inches(0.2), text, 7.5, False, COLORS["muted"])


def add_bullets(slide, items, left=0.75, top=1.25, width=5.8, height=4.8, size=15, color=None):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(7)
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color or COLORS["ink"]
    return box


def add_small_note(slide, text, left, top, width, height):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["light"]
    shape.line.color.rgb = COLORS["line"]
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Aptos"
    p.font.size = Pt(11)
    p.font.color.rgb = COLORS["ink"]
    p.space_after = Pt(0)
    return shape


def add_process(slide, steps, left=0.55, top=1.6, width=12.25):
    n = len(steps)
    box_w = width / n - 0.1
    for i, (label, caption, color) in enumerate(steps):
        x = left + i * (width / n)
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(top),
            Inches(box_w),
            Inches(1.15),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = color
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = label
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Aptos"
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = COLORS["white"]
        q = tf.add_paragraph()
        q.text = caption
        q.alignment = PP_ALIGN.CENTER
        q.font.name = "Aptos"
        q.font.size = Pt(8.5)
        q.font.color.rgb = COLORS["white"]
        if i < n - 1:
            add_textbox(slide, Inches(x + box_w + 0.03), Inches(top + 0.43), Inches(0.18), Inches(0.25), ">", 16, True, COLORS["muted"], PP_ALIGN.CENTER)


def add_table(slide, rows, left, top, widths, row_h=0.35, font_size=10):
    table = slide.shapes.add_table(len(rows), len(rows[0]), Inches(left), Inches(top), Inches(sum(widths)), Inches(row_h * len(rows))).table
    for col, width in enumerate(widths):
        table.columns[col].width = Inches(width)
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(value)
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLORS["light"] if r == 0 else COLORS["white"]
            for p in cell.text_frame.paragraphs:
                p.font.name = "Aptos"
                p.font.size = Pt(font_size)
                p.font.bold = r == 0
                p.font.color.rgb = COLORS["ink"]
    return table


def read_results():
    path = ROOT / "pet_cosmors_hpc_package" / "pet_solubility_results.csv"
    with path.open(newline="", encoding="utf-8") as stream:
        return list(csv.DictReader(stream))


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def slide(title, subtitle=None):
        s = prs.slides.add_slide(blank)
        add_title(s, title, subtitle)
        return s

    s = slide("PET solubility workflow with ORCA and openCOSMO-RS", "From official openCOSMO-RS examples to PET-in-GVL/formic-acid calculations")
    add_process(s, [
        ("Structure", "PET trimer + solvents", COLORS["blue"]),
        ("ORCA", "optimization + COSMORS SP", COLORS["green"]),
        ("openCOSMO-RS", "ln(gamma)", COLORS["amber"]),
        ("SLE", "x and g/g solvent", COLORS["red"]),
    ], top=2.0)
    add_small_note(s, "Final numerical results: GVL = 0.334 g PET/g solvent (infinite dilution); formic acid = 4.59e-4 g PET/g solvent.", 1.1, 4.2, 11.1, 0.9)
    add_footer(s, "Files: pet_cosmors_hpc_package/pet_solubility_results.csv; generated by orca_full_workflow/08_calculate_pet_solubility.py")

    s = slide("What openCOSMO-RS_py contributes", "Thermodynamic post-processing from COSMO surface information")
    add_bullets(s, [
        "openCOSMO-RS_py is a Python implementation of COSMO-RS with ORCA .orcacosmo parser support.",
        "The working API is COSMORS(par=openCOSMORS24a()), add_molecule([...]), add_job(x, T, refst='pure_component'), calculate().",
        "The model returns logarithmic activity coefficients ln(gamma); the solid solubility equation is a downstream thermodynamic calculation.",
        "Official examples include activity coefficients, sigma-profile parsing, and paracetamol solubility."
    ], width=11.7)
    add_footer(s, "Basis: openCOSMO-RS_py/README.md; openCOSMO-RS_py/examples/cosmors_calculation.ipynb cells 15, 27, 29")

    s = slide("OPI nightly documentation: role of ORCA", "ORCA creates the quantum-chemical COSMO input files needed by openCOSMO-RS")
    add_bullets(s, [
        "OPI tutorial workflow: prepare structures, run ORCA with COSMORS keyword, parse .orcacosmo, then compute solubility.",
        "Example ORCA setup uses calc.input.add_arbitrary_string('!COSMORS(ethanol)') and calc.input.ncores.",
        "The tutorial explicitly checks normal ORCA termination before parsing the .orcacosmo output.",
        "Our HPC implementation used direct ORCA input files instead of OPI, but kept the same COSMORS output target."
    ], width=11.7)
    add_footer(s, "Basis: openCOSMO-RS - OPI nightly Docs.pdf, pp. 1, 2, 4, 5")

    s = slide("Official solubility equation in the tutorial", "openCOSMO-RS gives ln(gamma); the SLE equation converts it to x")
    add_textbox(s, Inches(0.8), Inches(1.35), Inches(11.8), Inches(0.55), "ln(x) + ln(gamma) = -DeltaHfus/R * (1/T - 1/Tfus)", 26, True, COLORS["blue"], PP_ALIGN.CENTER)
    add_bullets(s, [
        "Non-iterative: evaluate ln(gamma) at infinite dilution, x = exp(rhs - ln(gamma_inf)).",
        "Iterative: solve ln(gamma(x)) + ln(x) - rhs = 0.",
        "Reference state used in code: refst='pure_component'.",
        "For literature-style screening, infinite-dilution values are the most directly comparable."
    ], top=2.35, width=11.7)
    add_footer(s, "Basis: openCOSMO-RS - OPI nightly Docs.pdf, pp. 10-12; openCOSMO-RS_py/examples/cosmors_calculation.ipynb cell 29")

    s = slide("Lee et al. 2025: PET solubility screening context", "The target paper uses COSMO-RS to screen dye/PET solubility across solvents")
    add_bullets(s, [
        "The article predicts PET and dye solubilities at 25 C and 100 C for 22 solvents.",
        "GVL and acetic/formic acid family solvents appear in the solvent-screening context; acetic acid is highlighted experimentally.",
        "The paper reports PET solubility in g solute per g solvent and uses the values to evaluate fibre preservation.",
        "COSMO-RS is used as an equilibrium thermodynamic screen; kinetic and transport effects are outside the model."
    ], width=11.7)
    add_footer(s, "Basis: Lee et al. 2025, Nature Sustainability, Fig. 1 and Methods, pp. 2-3, 8")

    s = slide("Lee et al. Methods: PET model and parameters", "These are the parameters we mirrored in the PET calculation")
    rows = [
        ["Item", "Lee et al.", "Our calculation"],
        ["PET model", "PET trimer, ends capped with EG motifs", "PET trimer from MD/optimized ORCA geometry"],
        ["MD selection", "20 PET conformers from one 10 ns MD trajectory", "1 conformer used in this first test"],
        ["Melting point", "260 C", "533.15 K"],
        ["Fusion enthalpy", "54.3 J/g", "converted to 34.674 kJ/mol using C32H30O14"],
        ["Output unit", "g/g solvent", "mole fraction + g/g solvent"],
    ]
    add_table(s, rows, 0.6, 1.3, [2.2, 5.1, 5.0], row_h=0.47, font_size=10)
    add_footer(s, "Basis: Lee et al. 2025 Methods, p. 8; our PET formula from pet_trimer_opt_final.xyz")

    s = slide("How our workflow maps to the literature workflow", "Same concept, different quantum-chemistry implementation")
    rows = [
        ["Stage", "Lee et al.", "This work"],
        ["Conformations", "20 PET conformers from MD", "single PET trimer conformer"],
        ["QM software", "Gaussian 16", "ORCA 6.0.0 on HPC"],
        ["COSMO-RS engine", "COSMOtherm BP_TZVP_22", "openCOSMO-RS_py openCOSMORS24a"],
        ["COSMO files", "COSMO files for PET/solvents/dyes", ".solute.orcacosmo for PET/GVL/formic acid"],
        ["Solubility equation", "chemical-potential based SLE", "equivalent ln(gamma)-based SLE form"],
    ]
    add_table(s, rows, 0.55, 1.25, [2.0, 5.2, 5.45], row_h=0.48, font_size=9.5)
    add_footer(s, "Basis: Lee et al. 2025 Methods, p. 8; openCOSMO-RS - OPI nightly Docs.pdf, pp. 10-12")

    s = slide("Local calculation files: geometry optimization", "Three optimized species were used for the COSMORS single-point step")
    rows = [
        ["Species", "Input/trajectory", "ORCA opt level", "Status"],
        ["Formic acid", "31-CH2O2/31-CH2O2_opt_trj.xyz", "M062X def2-TZVP D3zero Opt pal4", "converged"],
        ["GVL", "34-C5H8O2/34-C5H8O2_opt_trj.xyz", "M062X def2-TZVP D3zero Opt pal4", "converged"],
        ["PET trimer", "35-PET/35-PET_opt_trj.xyz", "M062X def2-SVP D3zero Opt pal64", "converged"],
    ]
    add_table(s, rows, 0.5, 1.3, [1.6, 4.5, 4.4, 1.7], row_h=0.55, font_size=9.5)
    add_footer(s, "Basis: 31-CH2O2/*_opt.out, 34-C5H8O2/*_opt.out, 35-PET/*_opt.out; files generated by 05/06 scripts")

    s = slide("Local calculation files: COSMORS single points", "Final optimized structures were converted to .orcacosmo files")
    rows = [
        ["Species", "ORCA input", "Output used in openCOSMO-RS", "Cores"],
        ["PET trimer", "pet_cosmors_hpc_package/pet_trimer.inp", "pet_trimer.solute.orcacosmo", "64"],
        ["GVL", "pet_cosmors_hpc_package/gvl.inp", "gvl.solute.orcacosmo", "4"],
        ["Formic acid", "pet_cosmors_hpc_package/formic_acid.inp", "formic_acid.solute.orcacosmo", "4"],
    ]
    add_table(s, rows, 0.5, 1.3, [1.7, 4.1, 4.8, 1.2], row_h=0.55, font_size=9.5)
    add_bullets(s, [
        "ORCA line used: BLYP def2-SVP noautostart miniprint palN COSMORS(ethanol).",
        "Only *.solute.orcacosmo files are used for PET-in-solvent calculations."
    ], left=0.75, top=4.25, width=11.4, height=1.2, size=13)
    add_footer(s, "Basis: orca_full_workflow/07_prepare_pet_cosmors_package.py; pet_cosmors_hpc_package/*.inp")

    s = slide("Post-processing script", "PET-specific wrapper around official openCOSMO-RS usage")
    add_bullets(s, [
        "Script: orca_full_workflow/08_calculate_pet_solubility.py",
        "Reads PET, GVL and formic-acid .solute.orcacosmo files.",
        "Uses openCOSMORS24a and refst='pure_component' to calculate ln(gamma).",
        "Converts DeltaHfus = 54.3 J/g to J/mol using molar mass from PET XYZ.",
        "Reports both mole fraction and g PET/g solvent."
    ], width=11.7)
    add_footer(s, "Basis: orca_full_workflow/08_calculate_pet_solubility.py")

    s = slide("Unit conversion used for PET", "Lee reports fusion enthalpy per gram; the SLE equation needs per mole")
    add_textbox(s, Inches(0.8), Inches(1.3), Inches(11.8), Inches(0.48), "PET trimer formula: C32H30O14; M = 638.5722 g/mol", 22, True, COLORS["blue"], PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.8), Inches(2.25), Inches(11.8), Inches(0.48), "DeltaHfus = 54.3 J/g * 638.5722 g/mol = 34.674 kJ/mol", 22, True, COLORS["green"], PP_ALIGN.CENTER)
    add_bullets(s, [
        "This conversion is required because the solid-liquid equilibrium equation uses molar thermodynamic quantities.",
        "For polymer models, the molar mass must match the molecular representation used in the COSMO-RS calculation.",
        "If a different PET oligomer length is used, the conversion must be updated."
    ], top=3.3, width=11.7, size=14)
    add_footer(s, "Basis: Lee et al. 2025 Methods p. 8; pet_cosmors_hpc_package/pet_trimer_opt_final.xyz; 08_calculate_pet_solubility.py")

    s = slide("Final numerical results", "Single-conformer PET trimer calculation at 298.15 K")
    data = read_results()
    rows = [["Solvent", "ln(gamma_inf)", "x_inf", "g/g_inf", "x_iter", "g/g_iter"]]
    for row in data:
        rows.append([
            "Formic acid" if row["solvent"] == "Formic_acid" else row["solvent"],
            "{:.4g}".format(float(row["ln_gamma_inf"])),
            "{:.4g}".format(float(row["x_inf"])),
            "{:.4g}".format(float(row["g_per_g_inf"])),
            "{:.4g}".format(float(row["x_iter"])),
            "{:.4g}".format(float(row["g_per_g_iter"])),
        ])
    add_table(s, rows, 0.65, 1.35, [1.8, 1.8, 1.6, 1.8, 1.7, 1.8], row_h=0.58, font_size=12)
    add_small_note(s, "Main interpretation: the calculated GVL value is much larger than formic acid. Using infinite-dilution g/g values, GVL is about 700x higher.", 1.0, 4.0, 11.1, 0.8)
    add_footer(s, "Basis: pet_cosmors_hpc_package/pet_solubility_results.csv")

    s = slide("Interpreting non-iterative vs iterative values", "Which number should be used for comparison?")
    add_bullets(s, [
        "x_inf and g/g_inf use the infinite-dilution activity coefficient.",
        "x_iter and g/g_iter solve the full composition-dependent equilibrium condition.",
        "Lee et al. describe chemical potentials at infinite dilution in the solvent; for screening comparison, g/g_inf is the cleanest first comparison.",
        "Iterative values are useful as a self-consistent correction when solubility is not extremely dilute."
    ], width=11.7)
    add_footer(s, "Basis: openCOSMO-RS - OPI nightly Docs.pdf, pp. 8, 10-12; Lee et al. 2025 Methods p. 8")

    s = slide("Important caveats", "Current result is a reproducibility test, not a full reproduction of Lee et al.")
    add_bullets(s, [
        "Only one PET conformation was used; Lee et al. used 20 PET trimer conformers from a 10 ns trajectory.",
        "ORCA/openCOSMO-RS parameterization differs from Gaussian/COSMOtherm BP_TZVP_22 in the paper.",
        "Our COSMORS single point used BLYP def2-SVP; higher basis or exact openCOSMO-RS recommended protocols may change values.",
        "COSMO-RS predicts equilibrium thermodynamics; degradation, diffusion, swelling and fibre morphology are not captured."
    ], width=11.7)
    add_footer(s, "Basis: Lee et al. 2025 Methods p. 8; Muller et al. 2025 openCOSMO-RS 24a discussion, p. 4")

    s = slide("Recommended next steps", "How to move from first test to paper-level comparison")
    add_bullets(s, [
        "Extract multiple PET conformers from MD, preferably using the same selection criterion: lowest internal potential energy.",
        "Run geometry optimization and COSMORS single-point calculations for each conformer.",
        "Generate conformer-weighted or ensemble-averaged ln(gamma)/chemical-potential estimates.",
        "Repeat for all target solvents and compare g/g_inf at 25 C and 100 C.",
        "Document deviations from Lee et al.: software, parameterization, basis set, solvent conformer handling."
    ], width=11.7)
    add_footer(s, "Basis: Lee et al. 2025 Methods p. 8; local scripts 06-08")

    s = slide("Appendix: local command sequence", "Commands used in this reproduction")
    add_bullets(s, [
        "python orca_full_workflow/07_prepare_pet_cosmors_package.py",
        "bash submit_all_cosmors.sh",
        "bash check_cosmors_results.sh",
        "python orca_full_workflow/08_calculate_pet_solubility.py --csv pet_cosmors_hpc_package/pet_solubility_results.csv",
        "Inputs: pet_trimer.solute.orcacosmo, gvl.solute.orcacosmo, formic_acid.solute.orcacosmo"
    ], width=11.7, size=13)
    add_footer(s, "Basis: pet_cosmors_hpc_package/README_submit.md; orca_full_workflow/08_calculate_pet_solubility.py")

    s = slide("Appendix: source map", "Where each claim in this deck comes from")
    rows = [
        ["Claim", "Source"],
        ["openCOSMO-RS API and examples", "openCOSMO-RS_py README + examples/cosmors_calculation.ipynb"],
        ["OPI/ORCA .orcacosmo workflow", "openCOSMO-RS - OPI nightly Docs.pdf, pp. 1-5"],
        ["Solubility formula and iterative/non-iterative methods", "openCOSMO-RS - OPI nightly Docs.pdf, pp. 10-12"],
        ["PET trimer, MD selection, 20 conformers", "Lee et al. 2025 Methods, p. 8"],
        ["PET Tm and DeltaHfus", "Lee et al. 2025 Methods, p. 8"],
        ["Our numerical results", "pet_cosmors_hpc_package/pet_solubility_results.csv"],
    ]
    add_table(s, rows, 0.5, 1.25, [4.3, 8.0], row_h=0.45, font_size=9.3)
    add_footer(s, "Generated deck: PET_openCOSMO_RS_workflow_summary.pptx")

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
