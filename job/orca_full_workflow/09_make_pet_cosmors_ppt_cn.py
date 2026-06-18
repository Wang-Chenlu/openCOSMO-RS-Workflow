#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""生成 PET/openCOSMO-RS 工作流中文汇报 PPT。"""

from pathlib import Path
import csv

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "PET_openCOSMO_RS_workflow_summary_CN.pptx"

FONT = "Microsoft YaHei"
COLORS = {
    "ink": RGBColor(31, 42, 55),
    "muted": RGBColor(91, 105, 122),
    "blue": RGBColor(37, 99, 235),
    "green": RGBColor(5, 150, 105),
    "red": RGBColor(190, 18, 60),
    "amber": RGBColor(180, 83, 9),
    "light": RGBColor(243, 246, 250),
    "line": RGBColor(205, 213, 223),
    "white": RGBColor(255, 255, 255),
}


def set_run(run, size=14, bold=False, color=None):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or COLORS["ink"]


def add_text(slide, x, y, w, h, text, size=14, bold=False, color=None, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    if align:
        p.alignment = align
    if p.runs:
        set_run(p.runs[0], size=size, bold=bold, color=color)
    return box


def add_title(slide, title, subtitle=None):
    add_text(slide, 0.55, 0.28, 12.2, 0.52, title, size=23, bold=True)
    if subtitle:
        add_text(slide, 0.58, 0.84, 12.0, 0.34, subtitle, size=10, color=COLORS["muted"])


def add_footer(slide, text):
    add_text(slide, 0.55, 7.13, 12.25, 0.24, text, size=7.5, color=COLORS["muted"])


def add_bullets(slide, items, x=0.72, y=1.28, w=11.85, h=5.35, size=14):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.space_after = Pt(8)
        p.level = 0
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = COLORS["ink"]
    return box


def add_note(slide, text, x, y, w, h, fill="light", size=12):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS[fill]
    shape.line.color.rgb = COLORS["line"]
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = FONT
    p.font.size = Pt(size)
    p.font.color.rgb = COLORS["ink"]
    return shape


def add_process(slide, steps, x=0.58, y=1.65, w=12.15):
    n = len(steps)
    step_w = w / n - 0.12
    for i, (title, caption, color) in enumerate(steps):
        left = x + i * (w / n)
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(y),
            Inches(step_w),
            Inches(1.18),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = color
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = title
        p.alignment = PP_ALIGN.CENTER
        p.font.name = FONT
        p.font.size = Pt(14.5)
        p.font.bold = True
        p.font.color.rgb = COLORS["white"]
        q = tf.add_paragraph()
        q.text = caption
        q.alignment = PP_ALIGN.CENTER
        q.font.name = FONT
        q.font.size = Pt(8.2)
        q.font.color.rgb = COLORS["white"]
        if i < n - 1:
            add_text(slide, left + step_w + 0.03, y + 0.43, 0.18, 0.25, ">", 15, True, COLORS["muted"], PP_ALIGN.CENTER)


def add_table(slide, rows, x, y, widths, row_h=0.42, size=9.5):
    table = slide.shapes.add_table(
        len(rows),
        len(rows[0]),
        Inches(x),
        Inches(y),
        Inches(sum(widths)),
        Inches(row_h * len(rows)),
    ).table
    for c, width in enumerate(widths):
        table.columns[c].width = Inches(width)
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(value)
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLORS["light"] if r == 0 else COLORS["white"]
            for p in cell.text_frame.paragraphs:
                p.font.name = FONT
                p.font.size = Pt(size)
                p.font.bold = r == 0
                p.font.color.rgb = COLORS["ink"]
    return table


def read_results():
    path = ROOT / "pet_cosmors_hpc_package" / "pet_solubility_results.csv"
    with path.open(newline="", encoding="utf-8") as stream:
        return list(csv.DictReader(stream))


def new_slide(prs, title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, title, subtitle)
    return slide


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    s = new_slide(prs, "PET 在 GVL/甲酸中的溶解度计算流程", "从 openCOSMO-RS_py 官方示例、OPI 文档、Lee 等 2025 论文到本次 ORCA/openCOSMO-RS 计算")
    add_process(s, [
        ("结构来源", "PET 三聚体 + GVL + 甲酸", COLORS["blue"]),
        ("ORCA", "几何优化 + COSMORS 单点", COLORS["green"]),
        ("openCOSMO-RS", "由 .orcacosmo 计算 ln(gamma)", COLORS["amber"]),
        ("固液平衡", "由 ln(gamma) 得到 x 与 g/g", COLORS["red"]),
    ], y=2.0)
    add_note(s, "本次计算结果：GVL = 0.334 g PET/g solvent；甲酸 = 4.59e-4 g PET/g solvent。这里采用无限稀释值 g/g_inf 作为与文献筛选结果的优先比较量。", 1.0, 4.15, 11.35, 0.9)
    add_footer(s, "结果文件：pet_cosmors_hpc_package/pet_solubility_results.csv；计算脚本：orca_full_workflow/08_calculate_pet_solubility.py")

    s = new_slide(prs, "1. openCOSMO-RS_py 的作用", "它不是量化软件，而是 COSMO-RS 热力学后处理程序")
    add_bullets(s, [
        "输入：量化计算得到的 COSMO 表面文件，例如 ORCA 生成的 .orcacosmo。",
        "核心 API：COSMORS(par=openCOSMORS24a())、add_molecule、add_job、calculate。",
        "输出：活度系数 ln(gamma)、sigma profile、面积/体积/氢键矩等描述符。",
        "溶解度本身不是直接从 ORCA 得到，而是由 openCOSMO-RS 的 ln(gamma) 结合固液平衡方程计算。",
    ])
    add_footer(s, "依据：openCOSMO-RS_py/README.md；openCOSMO-RS_py/examples/cosmors_calculation.ipynb")

    s = new_slide(prs, "2. OPI nightly 文档中的官方流程", "OPI 文档展示了 ORCA + openCOSMO-RS_py 的完整示例")
    add_bullets(s, [
        "文档目标：用 ORCA Python Interface 运行 ORCA，生成 .orcacosmo，并用 openCOSMO-RS_py 计算 sigma profile 和溶解度。",
        "关键 ORCA 关键词：calc.input.add_arbitrary_string('!COSMORS(ethanol)')。",
        "文档要求检查 ORCA 是否正常终止，然后再解析 .orcacosmo。",
        "我们在超算上不用 OPI，而是直接写 ORCA input；本质上复现的是同一个 COSMORS 单点输出流程。",
    ])
    add_footer(s, "依据：openCOSMO-RS - OPI nightly Docs.pdf，第 1、2、4、5 页")

    s = new_slide(prs, "3. 本次采用的溶解度公式", "openCOSMO-RS 先算 ln(gamma)，再通过固液平衡方程得到 x")
    add_text(s, 0.75, 1.35, 11.9, 0.55, "ln(x) + ln(gamma) = - ΔH_fus / R × (1/T - 1/T_fus)", 24, True, COLORS["blue"], PP_ALIGN.CENTER)
    add_bullets(s, [
        "非迭代：在无限稀释下计算 ln(gamma_inf)，x_inf = exp(rhs - ln(gamma_inf))。",
        "迭代：求解 ln(gamma(x)) + ln(x) - rhs = 0。",
        "openCOSMO-RS 调用中使用 refst='pure_component'，与官方示例一致。",
        "本式未加入热容修正项 Delta Cp，因此属于简化固液平衡模型。",
    ], y=2.35)
    add_footer(s, "依据：openCOSMO-RS - OPI nightly Docs.pdf，第 10-12 页；openCOSMO-RS_py/examples/cosmors_calculation.ipynb 中 paracetamol 示例")

    s = new_slide(prs, "4. Lee 等 2025 论文给出的目标体系", "论文用 COSMO-RS 筛选 PET 和染料在 22 种溶剂中的溶解度")
    add_bullets(s, [
        "论文比较 PET 与染料在不同溶剂中的溶解度，用于筛选“能溶染料但尽量不溶 PET”的萃取溶剂。",
        "Fig. 1 中给出 PET 和 Orange 30 在 25 °C 和 100 °C 的预测溶解度，单位为 g/g。",
        "GVL 在实验中可能导致染料转化；乙酸被认为在工艺上更适合回收和放大。",
        "我们这次先计算 GVL 和甲酸，目标是建立从 ORCA 到 openCOSMO-RS 的完整可跑通流程。",
    ])
    add_footer(s, "依据：Lee 等 2025, Nature Sustainability，第 2、3、7 页")

    s = new_slide(prs, "5. Lee 等 Methods 中的 PET 构型和参数", "这些参数决定了我们如何选择 PET 模型和热力学输入")
    rows = [
        ["项目", "Lee 等 2025", "本次计算"],
        ["PET 模型", "PET trimer，端基用 EG motif 封端", "PET 三聚体"],
        ["构型来源", "10 ns MD 轨迹中选 20 个低内能构象", "先使用 1 个 MD 构型测试"],
        ["熔点", "260 °C", "533.15 K"],
        ["熔融焓", "54.3 J/g", "用 PET 三聚体摩尔质量换算成 kJ/mol"],
        ["输出单位", "g/g solvent", "同时输出 mole fraction 和 g/g solvent"],
    ]
    add_table(s, rows, 0.55, 1.28, [2.0, 5.3, 5.0], row_h=0.48, size=9.3)
    add_footer(s, "依据：Lee 等 2025 Methods，第 8 页")

    s = new_slide(prs, "6. Lee 论文公式与本次公式的关系", "两者都在描述纯溶质与溶剂中无限稀释状态之间的化学势差")
    add_bullets(s, [
        "Lee 论文形式：log(x_i) = (mu_i^pure - mu_i^solvent - DeltaG_fus)/(RT ln 10)。",
        "openCOSMO-RS 示例形式：ln(x) + ln(gamma) = -DeltaH_fus/R × (1/T - 1/T_fus)。",
        "二者都需要纯组分参考态与溶剂中无限稀释状态的热力学差异。",
        "差异：Lee 使用 COSMOtherm/BP_TZVP_22；本次使用 openCOSMO-RS_py/openCOSMORS24a，并用 DeltaH_fus 近似融合自由能项。",
    ])
    add_footer(s, "依据：Lee 等 2025 Methods，第 8 页；OPI nightly Docs，第 10-12 页")

    s = new_slide(prs, "7. 我们实际生成和使用的文件", "从初始结构到最终 .orcacosmo")
    rows = [
        ["对象", "优化输入/输出", "COSMORS 输入", "openCOSMO-RS 使用文件"],
        ["PET", "35-PET/35-PET_opt.out", "pet_trimer.inp", "pet_trimer.solute.orcacosmo"],
        ["GVL", "34-C5H8O2/34-C5H8O2_opt.out", "gvl.inp", "gvl.solute.orcacosmo"],
        ["甲酸", "31-CH2O2/31-CH2O2_opt.out", "formic_acid.inp", "formic_acid.solute.orcacosmo"],
    ]
    add_table(s, rows, 0.45, 1.35, [1.25, 3.5, 2.75, 4.55], row_h=0.55, size=9.0)
    add_footer(s, "依据：31-CH2O2、34-C5H8O2、35-PET、pet_cosmors_hpc_package 目录")

    s = new_slide(prs, "8. ORCA 几何优化设置", "先优化构型，再做 COSMORS 单点")
    rows = [
        ["体系", "优化关键词", "核数", "终止状态"],
        ["甲酸", "M062X def2-TZVP D3zero Opt", "4", "正常收敛"],
        ["GVL", "M062X def2-TZVP D3zero Opt", "4", "正常收敛"],
        ["PET 三聚体", "M062X def2-SVP D3zero Opt", "64", "正常收敛"],
    ]
    add_table(s, rows, 0.7, 1.35, [2.0, 5.4, 1.4, 2.1], row_h=0.6, size=10)
    add_note(s, "PET 使用 def2-SVP 是为了在大体系上控制计算成本；小分子使用 def2-TZVP 保持较高几何精度。", 1.0, 4.55, 11.0, 0.65)
    add_footer(s, "依据：*_opt.inp 与 *_opt.out；检查项包括 THE OPTIMIZATION HAS CONVERGED 和 ORCA TERMINATED NORMALLY")

    s = new_slide(prs, "9. ORCA COSMORS 单点设置", "这一步生成 openCOSMO-RS 能读取的 .orcacosmo")
    add_text(s, 0.85, 1.35, 11.6, 0.45, "! BLYP def2-SVP noautostart miniprint palN COSMORS(ethanol)", 20, True, COLORS["blue"], PP_ALIGN.CENTER)
    add_bullets(s, [
        "PET: pal64；GVL 和甲酸: pal4。",
        "OPI 文档示例使用 !COSMORS(ethanol)；我们在直接 ORCA input 中保留这个关键词。",
        "对于自定义溶剂，真正用于 openCOSMO-RS 的是各自分子的 *.solute.orcacosmo。",
        "额外生成的 *.solvent.orcacosmo 是 ethanol 参考文件，本次 PET-in-GVL/甲酸计算不使用。",
    ], y=2.25)
    add_footer(s, "依据：openCOSMO-RS - OPI nightly Docs.pdf，第 4 页；pet_cosmors_hpc_package/*.inp")

    s = new_slide(prs, "10. PET 熔融焓单位换算", "论文给 J/g，溶解度公式需要 J/mol")
    add_text(s, 0.85, 1.25, 11.6, 0.45, "PET 三聚体公式：C32H30O14；M = 638.5722 g/mol", 22, True, COLORS["blue"], PP_ALIGN.CENTER)
    add_text(s, 0.85, 2.25, 11.6, 0.45, "DeltaH_fus = 54.3 J/g × 638.5722 g/mol = 34.6745 kJ/mol", 22, True, COLORS["green"], PP_ALIGN.CENTER)
    add_bullets(s, [
        "摩尔质量由 pet_trimer_opt_final.xyz 自动统计得到。",
        "如果后续改用不同聚合度的 PET 低聚物，必须重新换算。",
        "本次温度：T = 298.15 K；T_fus = 533.15 K。",
    ], y=3.35)
    add_footer(s, "依据：Lee 等 2025 Methods，第 8 页；orca_full_workflow/08_calculate_pet_solubility.py")

    s = new_slide(prs, "11. openCOSMO-RS 后处理脚本", "把 .orcacosmo 转换为 ln(gamma)、x 和 g/g")
    add_bullets(s, [
        "脚本：orca_full_workflow/08_calculate_pet_solubility.py。",
        "输入：pet_trimer.solute.orcacosmo、gvl.solute.orcacosmo、formic_acid.solute.orcacosmo。",
        "参数化：openCOSMORS24a。",
        "计算：先得到 ln(gamma_inf)，再计算 x_inf；同时求解迭代根 x_iter。",
        "单位转换：g/g = x_solute M_solute / [(1 - x_solute) M_solvent]。",
    ])
    add_footer(s, "依据：orca_full_workflow/08_calculate_pet_solubility.py；openCOSMO-RS_py 官方 API")

    s = new_slide(prs, "12. 最终计算结果", "单个 PET 三聚体构型，298.15 K")
    data = read_results()
    rows = [["溶剂", "ln(gamma_inf)", "x_inf", "g/g_inf", "x_iter", "g/g_iter"]]
    for row in data:
        solvent = "甲酸" if row["solvent"] == "Formic_acid" else row["solvent"]
        rows.append([
            solvent,
            "{:.4g}".format(float(row["ln_gamma_inf"])),
            "{:.4g}".format(float(row["x_inf"])),
            "{:.4g}".format(float(row["g_per_g_inf"])),
            "{:.4g}".format(float(row["x_iter"])),
            "{:.4g}".format(float(row["g_per_g_iter"])),
        ])
    add_table(s, rows, 0.65, 1.35, [1.6, 1.9, 1.6, 1.8, 1.7, 1.8], row_h=0.58, size=11)
    add_note(s, "结论：GVL 的预测 PET 溶解能力远高于甲酸。按 g/g_inf 比较，GVL 约为甲酸的 700 倍。", 1.0, 4.05, 11.15, 0.75)
    add_footer(s, "依据：pet_cosmors_hpc_package/pet_solubility_results.csv")

    s = new_slide(prs, "13. 应该比较哪一列？", "和文献筛选趋势对比时，优先看 g/g_inf")
    add_bullets(s, [
        "g/g_inf：基于无限稀释活度系数，和 Lee 论文中“溶质在溶剂中无限稀释化学势”的描述最接近。",
        "g/g_iter：考虑组成变化后的自洽修正；当溶解度不再极稀时有参考意义。",
        "本次 GVL 的 x_inf 约 0.05，迭代修正明显；甲酸极稀，两个值几乎相同。",
        "用于初筛排序：g/g_inf；用于更严格估计：同时报告 g/g_inf 与 g/g_iter。",
    ])
    add_footer(s, "依据：OPI nightly Docs，第 10-12 页；Lee 等 2025 Methods，第 8 页")

    s = new_slide(prs, "14. 与 Lee 等工作的差异和限制", "本次是可跑通流程和趋势测试，不是完全复现")
    add_bullets(s, [
        "构象数量：Lee 使用 20 个 PET 三聚体构象；本次使用 1 个构象。",
        "软件和参数化：Lee 使用 Gaussian + COSMOtherm BP_TZVP_22；本次使用 ORCA + openCOSMO-RS_py openCOSMORS24a。",
        "量化级别：本次 COSMORS 单点为 BLYP def2-SVP；更高基组可能改变定量结果。",
        "模型边界：COSMO-RS 是平衡热力学模型，不描述 PET 纤维形貌、扩散、降解和溶胀动力学。",
    ])
    add_footer(s, "依据：Lee 等 2025 Methods，第 8 页；Müller 等 2025 关于 openCOSMO-RS 24a 的讨论")

    s = new_slide(prs, "15. 下一步如何做到更接近论文", "从单构型测试推进到构象集合")
    add_bullets(s, [
        "从 MD 轨迹中按内能选择多个 PET 构象，目标接近 Lee 的 20 个构象。",
        "为每个 PET 构象进行 ORCA 优化和 COSMORS 单点计算。",
        "对 GVL、甲酸以及更多候选溶剂重复小分子构象搜索和 COSMORS 单点。",
        "进行构象加权或集合平均，再输出 g/g_inf 与 g/g_iter。",
        "同时在 298.15 K 和 373.15 K 计算，便于对应 25 °C 和 100 °C。"
    ])
    add_footer(s, "依据：Lee 等 2025 Methods，第 8 页；本地脚本 06-08 可扩展到多构象")

    s = new_slide(prs, "16. 本次可复现命令", "从提交包到结果表")
    add_bullets(s, [
        "python orca_full_workflow/07_prepare_pet_cosmors_package.py",
        "bash submit_all_cosmors.sh",
        "bash check_cosmors_results.sh",
        "python orca_full_workflow/08_calculate_pet_solubility.py --csv pet_cosmors_hpc_package/pet_solubility_results.csv",
        "最终输入文件：pet_trimer.solute.orcacosmo、gvl.solute.orcacosmo、formic_acid.solute.orcacosmo",
    ], size=13)
    add_footer(s, "依据：pet_cosmors_hpc_package/README_submit.md；orca_full_workflow/07_prepare_pet_cosmors_package.py；08_calculate_pet_solubility.py")

    s = new_slide(prs, "17. 依据索引", "PPT 中每类信息对应的来源")
    rows = [
        ["内容", "来源"],
        ["openCOSMO-RS_py 的 API 和示例", "openCOSMO-RS_py/README.md；examples/cosmors_calculation.ipynb"],
        ["ORCA/OPI 生成 .orcacosmo 的流程", "openCOSMO-RS - OPI nightly Docs.pdf，第 1-5 页"],
        ["非迭代/迭代溶解度公式", "openCOSMO-RS - OPI nightly Docs.pdf，第 10-12 页"],
        ["PET 三聚体、MD、20 构象", "Lee 等 2025 Methods，第 8 页"],
        ["PET 熔点与熔融焓", "Lee 等 2025 Methods，第 8 页"],
        ["本次 ORCA/openCOSMO-RS 文件和结果", "31-CH2O2、34-C5H8O2、35-PET、pet_cosmors_hpc_package"],
    ]
    add_table(s, rows, 0.55, 1.25, [3.9, 8.45], row_h=0.45, size=9.0)
    add_footer(s, "生成文件：PET_openCOSMO_RS_workflow_summary_CN.pptx")

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
