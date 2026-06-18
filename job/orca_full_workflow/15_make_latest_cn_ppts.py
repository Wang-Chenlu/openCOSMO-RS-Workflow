#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Generate Chinese PPTs for the latest PET results and openCOSMO-RS Book guide."""

from pathlib import Path
import csv
import math

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt


JOB = Path(__file__).resolve().parents[1]
ROOT = JOB.parent
OPEN_DIR = ROOT / "openCOSMO"

RESULTS_PPT = JOB / "PET_openCOSMO_RS_latest_results_CN.pptx"
BOOK_PPT = OPEN_DIR / "openCOSMO_RS_Book_Chinese_Manual.pptx"

FONT = "Microsoft YaHei"
COLORS = {
    "ink": RGBColor(28, 36, 48),
    "muted": RGBColor(89, 99, 114),
    "blue": RGBColor(35, 92, 168),
    "green": RGBColor(34, 122, 94),
    "red": RGBColor(172, 54, 54),
    "amber": RGBColor(176, 117, 36),
    "violet": RGBColor(109, 82, 161),
    "light": RGBColor(243, 246, 250),
    "line": RGBColor(205, 213, 223),
    "white": RGBColor(255, 255, 255),
}

SOURCE_BOOK = "https://priscillaensleysolis.github.io/openCOSMO_RS_Book/intro.html"


def prs_base():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def set_paragraph_font(paragraph, size=13, bold=False, color=None):
    paragraph.font.name = FONT
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color or COLORS["ink"]


def add_text(slide, x, y, w, h, text, size=13, bold=False, color=None, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.text = text
    if align is not None:
        p.alignment = align
    set_paragraph_font(p, size=size, bold=bold, color=color)
    return box


def add_title(slide, title, subtitle=None):
    add_text(slide, 0.55, 0.28, 12.15, 0.45, title, size=23, bold=True)
    if subtitle:
        add_text(slide, 0.58, 0.82, 12.1, 0.33, subtitle, size=9.5, color=COLORS["muted"])


def add_footer(slide, text):
    add_text(slide, 0.55, 7.12, 12.2, 0.24, text, size=7.2, color=COLORS["muted"])


def new_slide(prs, title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, title, subtitle)
    return slide


def add_bullets(slide, items, x=0.72, y=1.28, w=11.85, h=5.35, size=13):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.space_after = Pt(7)
        p.level = 0
        set_paragraph_font(p, size=size)
    return box


def add_note(slide, text, x, y, w, h, fill="light", size=12):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS[fill]
    shape.line.color.rgb = COLORS["line"]
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.text = text
    set_paragraph_font(p, size=size)
    return shape


def add_process(slide, steps, x=0.58, y=1.65, w=12.15):
    n = len(steps)
    step_w = w / n - 0.12
    for i, (title, caption, color_key) in enumerate(steps):
        left = x + i * (w / n)
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(y),
            Inches(step_w),
            Inches(1.15),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLORS[color_key]
        shape.line.color.rgb = COLORS[color_key]
        tf = shape.text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.alignment = PP_ALIGN.CENTER
        set_paragraph_font(p, size=14.5, bold=True, color=COLORS["white"])
        q = tf.add_paragraph()
        q.text = caption
        q.alignment = PP_ALIGN.CENTER
        set_paragraph_font(q, size=8.4, color=COLORS["white"])
        if i < n - 1:
            add_text(slide, left + step_w + 0.02, y + 0.43, 0.18, 0.25, ">", 15, True, COLORS["muted"], PP_ALIGN.CENTER)


def add_table(slide, rows, x, y, widths, row_h=0.42, size=9.5):
    table = slide.shapes.add_table(
        len(rows),
        len(rows[0]),
        Inches(x),
        Inches(y),
        Inches(sum(widths)),
        Inches(row_h * len(rows)),
    ).table
    for c, width in enumerate(widths):
        table.columns[c].width = Inches(width)
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(value)
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLORS["light"] if r == 0 else COLORS["white"]
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                set_paragraph_font(p, size=size, bold=(r == 0))
    return table


def add_bar_chart(slide, rows, x=0.75, y=1.55, w=11.6, h=4.5):
    max_value = max(row["mean_iter"] for row in rows)
    bar_h = h / len(rows) - 0.18
    for i, row in enumerate(rows):
        top = y + i * (h / len(rows))
        add_text(slide, x, top + 0.03, 1.25, 0.28, row["label"], size=10.5, bold=True)
        track = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(x + 1.5),
            Inches(top),
            Inches(w - 2.6),
            Inches(bar_h),
        )
        track.fill.solid()
        track.fill.fore_color.rgb = COLORS["light"]
        track.line.color.rgb = COLORS["line"]
        bar_w = (w - 2.6) * row["mean_iter"] / max_value
        bar = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(x + 1.5),
            Inches(top),
            Inches(bar_w),
            Inches(bar_h),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = COLORS[row["color"]]
        bar.line.color.rgb = COLORS[row["color"]]
        add_text(
            slide,
            x + 1.62 + min(bar_w, w - 3.2),
            top + 0.02,
            1.0,
            0.25,
            "{:.3g}".format(row["mean_iter"]),
            size=9,
            bold=True,
            color=COLORS["ink"],
        )


def read_csv(path):
    with path.open(newline="", encoding="utf-8") as stream:
        return list(csv.DictReader(stream))


def get_summary_rows():
    base_rows = read_csv(JOB / "pet_1_10ns_solubility_results" / "pet_1_10ns_solubility_summary.csv")
    extra_rows = read_csv(JOB / "pet_extra_solvent_solubility_results" / "pet_extra_solvent_solubility_summary.csv")
    source = {}
    for row in base_rows + extra_rows:
        source[(row["solvent"], row["quantity"])] = row

    ordered = [
        ("GVL", "GVL", "blue"),
        ("NMP", "NMP", "green"),
        ("Formic_acid", "CH2O2 / 甲酸", "amber"),
        ("Isopropanol", "ISOPROPANOL", "violet"),
    ]
    rows = []
    for solvent, label, color in ordered:
        inf = source[(solvent, "g_per_g_inf")]
        itr = source[(solvent, "g_per_g_iter")]
        rows.append(
            {
                "solvent": solvent,
                "label": label,
                "color": color,
                "mean_inf": float(inf["mean"]),
                "median_inf": float(inf["median"]),
                "stdev_inf": float(inf["stdev"]),
                "min_inf": float(inf["min"]),
                "max_inf": float(inf["max"]),
                "n_nonfinite_inf": int(inf.get("n_nonfinite") or 0),
                "mean_iter": float(itr["mean"]),
                "median_iter": float(itr["median"]),
                "stdev_iter": float(itr["stdev"]),
                "min_iter": float(itr["min"]),
                "max_iter": float(itr["max"]),
            }
        )
    return rows


def fmt(value):
    if isinstance(value, float) and (math.isnan(value) or math.isinf(value)):
        return str(value)
    return "{:.6g}".format(value)


def make_results_ppt():
    rows = get_summary_rows()
    prs = prs_base()

    s = new_slide(
        prs,
        "PET 在四种溶剂中的 openCOSMO-RS 溶解度预测",
        "最新结果汇总：GVL、NMP、CH2O2/甲酸、ISOPROPANOL；10 个 PET 构型，T = 298.15 K",
    )
    add_process(
        s,
        [
            ("MD 构型", "1-10 ns PET 三聚体", "blue"),
            ("ORCA", "优化 + COSMORS 单点", "green"),
            ("openCOSMO-RS", ".orcacosmo -> ln(gamma)", "amber"),
            ("溶解度", "x 与 g/g solvent", "red"),
        ],
        y=1.65,
    )
    add_note(
        s,
        "建议主表采用 g/g_iter：NMP 的无限稀释近似在 5/10 个构型中给出 x_inf_raw > 1，导致 g/g_inf 为 inf，不能直接平均。",
        0.9,
        3.65,
        11.5,
        0.75,
    )
    add_footer(s, "数据来源：job/pet_1_10ns_solubility_results 与 job/pet_extra_solvent_solubility_results")

    s = new_slide(prs, "核心结论", "按平均 g/g_iter 排序：NMP > GVL > ISOPROPANOL > CH2O2/甲酸")
    add_bar_chart(s, rows, y=1.45)
    add_note(
        s,
        "NMP 预测溶解能力最高；GVL 次之；异丙醇明显低于 GVL；甲酸在本模型下对 PET 的溶解度最低。",
        1.0,
        6.05,
        11.15,
        0.62,
    )
    add_footer(s, "指标：10 个 PET 构型的 g PET / g solvent，采用迭代解 g/g_iter")

    s = new_slide(prs, "四种溶剂汇总表", "单位：g PET / g solvent；按用户指定顺序排列")
    table_rows = [["溶剂", "mean", "median", "stdev", "min", "max"]]
    for row in rows:
        table_rows.append(
            [
                row["label"],
                fmt(row["mean_iter"]),
                fmt(row["median_iter"]),
                fmt(row["stdev_iter"]),
                fmt(row["min_iter"]),
                fmt(row["max_iter"]),
            ]
        )
    add_table(s, table_rows, 0.75, 1.3, [2.3, 1.65, 1.65, 1.65, 1.65, 1.65], row_h=0.55, size=10.5)
    add_footer(s, "采用 g/g_iter；T = 298.15 K；PET: DeltaHfus = 54.3 J/g, Tfus = 533.15 K")

    s = new_slide(prs, "为什么 NMP 不建议用 g/g_inf 平均值", "无限稀释近似可以筛选趋势，但在高溶解度体系可能越界")
    table_rows = [["溶剂", "g/g_inf mean", "g/g_iter mean", "g/g_inf 非有限值", "建议"]]
    for row in rows:
        table_rows.append(
            [
                row["label"],
                fmt(row["mean_inf"]),
                fmt(row["mean_iter"]),
                "{}/10".format(row["n_nonfinite_inf"]),
                "看 iter" if row["n_nonfinite_inf"] else "inf/iter 均可报告",
            ]
        )
    add_table(s, table_rows, 0.55, 1.28, [2.0, 2.0, 2.0, 2.0, 3.0], row_h=0.52, size=9.5)
    add_note(s, "当 x_inf_raw >= 1 时，x_solvent = 1 - x_solute 接近 0，g/g 的换算会发散。因此 NMP 的非迭代 g/g_inf 不应作为平均比较值。", 0.9, 4.55, 11.35, 0.75)
    add_footer(s, "脚本：job/orca_full_workflow/14_calculate_pet_extra_solvents_solubility.py")

    s = new_slide(prs, "计算公式与单位换算", "openCOSMO-RS 给出 ln(gamma)，溶解度由固液平衡方程得到")
    add_text(s, 0.85, 1.3, 11.6, 0.5, "ln(x) + ln(gamma) = - DeltaH_fus / R * (1/T - 1/T_fus)", 21, True, COLORS["blue"], PP_ALIGN.CENTER)
    add_bullets(
        s,
        [
            "非迭代：使用无限稀释 ln(gamma_inf)，x_inf = exp(rhs - ln(gamma_inf))。",
            "迭代：求解 ln(gamma(x)) + ln(x) - rhs = 0，得到 x_iter。",
            "质量比：g/g = x_solute * M_solute / [(1 - x_solute) * M_solvent]。",
            "本次 PET 三聚体摩尔质量：638.5722 g/mol；DeltaHfus = 54.3 J/g = 34.6745 kJ/mol。",
        ],
        y=2.25,
    )
    add_footer(s, "依据：openCOSMO-RS OPI Docs 溶解度示例；本地脚本 12 与 14")

    s = new_slide(prs, "文件和脚本索引", "整理后的目录结构中，正式计算都在 job 下")
    add_table(
        s,
        [
            ["内容", "路径"],
            ["GVL/甲酸汇总", "job/pet_1_10ns_solubility_results/*.csv"],
            ["NMP/异丙醇汇总", "job/pet_extra_solvent_solubility_results/*.csv"],
            ["PET 02-10 ns COSMO 文件", "job/pet_02_10ns_cosmors_package/*.orcacosmo"],
            ["小分子 COSMO 文件", "job/pet_cosmors_hpc_package 与 job/pet_extra_solvent_cosmors_package"],
            ["重算 GVL/甲酸", "python job/orca_full_workflow/12_calculate_pet_1_10ns_solubility.py"],
            ["重算 NMP/异丙醇", "python job/orca_full_workflow/14_calculate_pet_extra_solvents_solubility.py"],
        ],
        0.55,
        1.25,
        [3.0, 8.9],
        row_h=0.48,
        size=9.2,
    )
    add_footer(s, "目录说明：README_folder_layout.md")

    prs.save(RESULTS_PPT)
    print(RESULTS_PPT)


def make_book_ppt():
    prs = prs_base()

    s = new_slide(
        prs,
        "openCOSMO-RS Book 中文手册",
        "基于 Priscilla Debanhi Ensley Solis、Simon Müller、Edgar Ivan Sanchez Medina 的在线指南整理",
    )
    add_note(
        s,
        "这份手册按“理论 -> 构象管线 -> 输入文件 -> 运行 -> .orcacosmo 内容 -> openCOSMO-RS API -> 结果分析”的顺序组织。",
        0.95,
        2.15,
        11.4,
        0.8,
    )
    add_footer(s, "来源：{}；网页说明其目标是 Windows 上 openCOSMO-RS 的逐步教程".format(SOURCE_BOOK))

    s = new_slide(prs, "1. 这本 Book 解决什么问题", "从安装、生成 .orcacosmo 到调用 openCOSMO-RS 分析结果")
    add_bullets(
        s,
        [
            "指南目标：介绍理论概念、软件需求、输入文件准备、示例应用和结果分析。",
            "应用目标：预测溶剂化现象和流体体系热力学性质。",
            "两部分结构：第一部分讲 COSMO-RS 与 conformer pipeline；第二部分讲 openCOSMO-RS 的 Python 调用。",
            "对本项目的意义：帮助理解 ORCA 生成 .orcacosmo、openCOSMO-RS 读取并计算 ln(gamma) 的全过程。",
        ],
    )
    add_footer(s, "来源：openCOSMO-RS Book intro.html")

    s = new_slide(prs, "2. COSMO-RS 的核心思想", "从分子结构预测液相热力学性质")
    add_bullets(
        s,
        [
            "COSMO-RS 使用分子表面上的屏蔽电荷密度来描述分子间相互作用。",
            "它可以估算分配系数、溶剂化自由能、活度系数等热力学量。",
            "量化软件先计算理想导体环境下的分子表面电荷；COSMO-RS 再把这些表面片段放入真实液体模型中。",
            "sigma profile 是表面电荷密度的分布图，类似分子表面极性特征的直方图。",
        ],
    )
    add_footer(s, "来源：Book I. What is COSMO-RS；COSMO-RS comes into play；The sigma profiles")

    s = new_slide(prs, "3. openCOSMO-RS 与 24a 参数化", "开源 COSMO-RS 实现，结合 ORCA 工作流")
    add_bullets(
        s,
        [
            "openCOSMO-RS 是开源模型，便于学术用户使用、评估和开发。",
            "模型不仅使用主屏蔽电荷密度 sigma，也使用额外片段描述符来提高预测能力。",
            "Book 强调 RDKit/ORCA 工作流：RDKit 生成构象，ORCA 完成量化计算。",
            "openCOSMO-RS 24a 是改进版本，面向使用 ORCA 量化计算的溶剂化自由能预测。",
        ],
    )
    add_footer(s, "来源：Book: What is openCOSMO-RS; openCOSMO-RS 24a")

    s = new_slide(prs, "4. conformer pipeline 需要什么", "它负责从分子结构出发生成 openCOSMO-RS 可用输入")
    add_table(
        s,
        [
            ["组件", "作用"],
            ["ORCA 6.0++", "量子化学计算；生成分子化学信息与性质"],
            ["xtb", "构象搜索和快速筛选"],
            ["Balloon", "可选；某些分子的构象预测会需要"],
            ["RDKit", "无 xyz 时从 SMILES 生成初始构象"],
        ],
        0.85,
        1.35,
        [2.2, 8.9],
        row_h=0.58,
        size=10.5,
    )
    add_footer(s, "来源：Book: 1. openCOSMO-RS_conformer_pipeline")

    s = new_slide(prs, "5. pipeline 输入文件格式", "每一行描述一个分子")
    add_text(s, 0.8, 1.35, 11.75, 0.5, "name [TAB] SMILES [TAB] xyz_file [TAB] charge [TAB] geometry_optimization", 18, True, COLORS["blue"], PP_ALIGN.CENTER)
    add_bullets(
        s,
        [
            "name：任意分子名，例如 water、methanol。",
            "SMILES：用字符串表示分子结构，例如水为 O，甲醇为 CO。",
            "xyz_file：可选；如果不提供，程序会用 RDKit 或 Balloon 生成构象，再交给 xtb 搜索。",
            "charge：中性分子为 0，阳离子可为 +1/+2，阴离子可为 -1/-2。",
            "geometry_optimization：TRUE 表示优化；FALSE 或空白表示不做优化。",
        ],
        y=2.25,
    )
    add_footer(s, "来源：Book: 4. Let's create the input file")

    s = new_slide(prs, "6. 如何调用 conformer pipeline", "命令行指定结构文件、CPCM 半径和核数")
    add_text(s, 0.9, 1.35, 11.5, 0.45, "python ConformerGenerator.py --structures_file file.inp --cpcm_radii cpcm_radii.inp --n_cores 2", 16, True, COLORS["green"], PP_ALIGN.CENTER)
    add_bullets(
        s,
        [
            "ConformerGenerator.py 是主程序。",
            "file.inp 是上一页创建的分子列表文件。",
            "cpcm_radii.inp 指定原子半径。",
            "n_cores 控制使用的 CPU 核数。",
            "程序结束时出现 La fin 表示运行完整结束。",
        ],
        y=2.25,
    )
    add_footer(s, "来源：Book: 5. Calling the openCOSMO-RS_conformer_pipeline program")

    s = new_slide(prs, "7. pipeline 计算过程如何理解", "主要分成气相计算和 CPCM 计算")
    add_bullets(
        s,
        [
            "Gas Calculation：考虑气相中的溶质，是计算溶剂化自由能所需的一部分。",
            "CPCM Calculation：把溶剂作为连续极化介质包围溶质，用于计算活度系数相关信息。",
            "构象数量由可旋转键数启发式决定：少于等于 7 个通常生成 50 个构象；8-12 个为 200 个；13 个以上为 300 个。",
            "流程包括能量排序、能量窗口过滤、RMSD 去重、快速 DFT、最终 DFT 以及单点能计算。",
        ],
    )
    add_footer(s, "来源：Book: 6. Analysis of the Calculations and the Input File")

    s = new_slide(prs, "8. .orcacosmo 文件里有什么", "它是 ORCA 到 openCOSMO-RS 的桥梁")
    add_table(
        s,
        [
            ["区块", "含义"],
            ["ENERGY", "CPCM/理想导体环境下得到的能量"],
            ["DIPOLE MOMENT", "偶极矩的 x/y/z 分量"],
            ["XYZ_FILE", "原子数与笛卡尔坐标"],
            ["COSMO", "表面点、介电信息、介电能等 COSMO 信息"],
            ["Surface points", "片段中心、面积、化学势、原子索引等"],
            ["COSMO_corrected", "离腔电荷修正后的介电能与 CPCM 电荷"],
        ],
        0.55,
        1.15,
        [2.35, 9.55],
        row_h=0.46,
        size=9.5,
    )
    add_footer(s, "来源：Book: 6.1 What is in our input file?")

    s = new_slide(prs, "9. openCOSMO-RS 的 Python 调用", "读取 .orcacosmo，加入分子、组成和温度")
    add_bullets(
        s,
        [
            "安装方式：pip install git+https://github.com/TUHH-TVT/openCOSMO-RS_py。",
            "常用导入：COSMORS、SigmaProfileParser、openCOSMORS24a、可视化函数。",
            "初始化：par = openCOSMORS24a(); model = COSMORS(par)。",
            "加入分子：model.add_molecule([path_to_orcacosmo])。",
            "加入任务：model.add_job(mole_fractions, T, refst)。",
            "计算：results = model.calculate()。",
        ],
    )
    add_footer(s, "来源：Book: II. openCOSMO-RS; First Examples")

    s = new_slide(prs, "10. refst 怎么选", "参考态决定 ln(gamma) 的物理含义")
    add_table(
        s,
        [
            ["refst", "用途"],
            ["pure_component", "相对纯组分；常用于无限稀释活度系数和混合物偏离纯态的比较"],
            ["reference_mixture", "相对某个指定参考混合物；用于比较两个组成状态"],
            ["cosmo", "不使用常规参考态；常用于溶剂化能相关计算"],
        ],
        0.75,
        1.35,
        [2.3, 9.4],
        row_h=0.62,
        size=10.2,
    )
    add_note(s, "本项目计算 PET 溶解度时使用 refst='pure_component'，并在 x_solute -> 0 时得到无限稀释 ln(gamma_inf)。", 1.0, 4.3, 11.2, 0.7)
    add_footer(s, "来源：Book: First Examples 2.4")

    s = new_slide(prs, "11. results 字典如何读", "openCOSMO-RS 输出总项、残余项和组合项")
    add_bullets(
        s,
        [
            "tot['lng']：总对数活度系数 ln(gamma)，等于残余项与组合项之和。",
            "enth['lng']：残余/焓贡献，反映片段间具体相互作用与参考态的差异。",
            "comb['lng']：组合项，来自分子大小、形状和混合熵差异。",
            "ln(gamma) > 0：比理想溶液相互作用更弱；ln(gamma) < 0：相互作用更有利。",
        ],
    )
    add_footer(s, "来源：Book: 3 Results Analysis")

    s = new_slide(prs, "12. 可视化结果怎么理解", "sigma profile 和 3D segment plot 是诊断分子表面的工具")
    add_bullets(
        s,
        [
            "plot_3D_segment_location：显示分子表面片段的 3D 位置；点大小与片段面积相关，颜色代表片段电荷。",
            "plot_sigma_profiles：x 轴为 sigma，即屏蔽电荷密度；y 轴为有效概率函数。",
            "sigma profile 表示不同电荷密度对应的表面积，可用于判断分子极性和相互作用潜力。",
            "extended sigma profile 增加正交表面电荷密度，用于观察更丰富的表面电荷分布。",
        ],
    )
    add_footer(s, "来源：Book: 3 Results Analysis")

    s = new_slide(prs, "13. 与本 PET 项目的对应关系", "我们没有直接运行 pipeline，但实现了同样的物理链条")
    add_table(
        s,
        [
            ["Book 流程", "本项目流程"],
            ["SMILES/xyz -> 构象搜索", "PET 构型来自 MD；小分子来自 xyz"],
            ["ORCA 量化计算", "超算上 ORCA 优化 + COSMORS 单点"],
            [".orcacosmo 输出", "pet_*.solute.orcacosmo 与 solvent.solute.orcacosmo"],
            ["COSMORS.add_molecule", "脚本 12 与 14 加入 PET 和溶剂文件"],
            ["results['tot']['lng']", "用于固液平衡方程计算 PET 溶解度"],
        ],
        0.7,
        1.25,
        [4.4, 7.4],
        row_h=0.52,
        size=9.5,
    )
    add_footer(s, "来源：Book 各章节；本地 job/orca_full_workflow 脚本")

    s = new_slide(prs, "14. 实操检查清单", "以后扩展新溶剂或新 PET 构型时按这张表检查")
    add_bullets(
        s,
        [
            "1. 分子结构明确：SMILES 或 xyz；聚合物低聚体构型来源要记录。",
            "2. ORCA 优化正常结束：检查 ORCA TERMINATED NORMALLY 和优化收敛。",
            "3. COSMORS 单点正常结束：确认生成 solute.orcacosmo。",
            "4. openCOSMO-RS 能读入文件：add_molecule 不报错。",
            "5. 结果解释时区分 ln(gamma_inf)、x_inf、x_iter、g/g_inf、g/g_iter。",
            "6. 若 x_inf_raw > 1，非迭代质量比会发散，优先报告迭代解。",
        ],
    )
    add_footer(s, "整理依据：openCOSMO-RS Book；本项目 NMP 结果中的 x_inf_raw > 1 情况")

    s = new_slide(prs, "15. 资料索引", "PPT 中用到的主要来源")
    add_table(
        s,
        [
            ["主题", "来源"],
            ["Book 入口与教程目标", SOURCE_BOOK],
            ["COSMO-RS 理论", "notebooks/1/opencosmorsintro.html"],
            ["conformer pipeline", "notebooks/1/conformer pipeline.html"],
            ["输入文件格式", "notebooks/4/Input file.html"],
            ["运行命令", "notebooks/5/Call program.html"],
            [".orcacosmo 内容", "notebooks/6/overview.html"],
            ["openCOSMO-RS API 和结果分析", "notebooks/7/firstexamples.html 与 ResultsAnalysis.html"],
        ],
        0.55,
        1.25,
        [3.3, 8.7],
        row_h=0.47,
        size=8.5,
    )
    add_footer(s, "生成文件：openCOSMO_RS_Book_Chinese_Manual.pptx")

    prs.save(BOOK_PPT)
    print(BOOK_PPT)


def main():
    OPEN_DIR.mkdir(parents=True, exist_ok=True)
    make_results_ppt()
    make_book_ppt()
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
